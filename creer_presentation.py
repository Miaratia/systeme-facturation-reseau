from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# Palette de couleurs
NAVY = RGBColor(16, 44, 87)
GRAY = RGBColor(100, 100, 100)

slides_data = [
    {
        "title": "Conception et Réalisation d'un Système de Facturation Réseau Convergent",
        "subtitle": "Traitement en temps réel des événements réseau (CDRs) dans un système réparti\n\nPrésenté par : Larissa\nUPH - Maintenance & Systèmes Réseaux | ADS 360",
        "is_title_slide": True
    },
    {
        "title": "Partie 1 : Contexte Réseau & Problématique",
        "bullets": [
            "Contexte Télécom : Évolution vers la convergence des services (Voix, SMS, Data) sur réseau IP.",
            "Capture de données brutes : Génération continue de CDRs (Call Detail Records).",
            "Problématique : Comment automatiser le débit en temps réel des comptes abonnés selon une grille tarifaire dynamique ?"
        ]
    },
    {
        "title": "Partie 1 : Architecture du Système Réparti",
        "bullets": [
            "Domaine : Télécommunications & Téléinformatique.",
            "Modèle Architectural : Client-Serveur / N-Tiers.",
            "Couche Collecte : Équipements réseau générant les CDRs.",
            "Couche Métier (PHP) : Moteur de tarification (Rating Engine) et règles de gestion.",
            "Couche Données (MySQL) : Gestion des comptes abonnés et persistance."
        ]
    },
    {
        "title": "Partie 2 : Modélisation de la Base de Données",
        "bullets": [
            "Table 'abonnes' : Stocke le nom, le type de compte, le solde argent (MGA) et le solde Data (Mo).",
            "Table 'tarifs' : Définition des prix unitaires par service (VOIX, SMS, DATA).",
            "Table 'cdrs' : Registre de tous les événements réseau (Statuts : FACTURE ou REJETE).",
            "Table 'consommations_facturees' : Traçabilité et historique financier des débits."
        ]
    },
    {
        "title": "Partie 2 : Le Moteur de Tarification (Rating Engine)",
        "bullets": [
            "Priorité Data : Déduction prioritaire du solde de Mo gratuits. Facturation en argent si le forfait est épuisé.",
            "Contrôle de Solde : Vérification de la solvabilité du compte prépayé avant autorisation.",
            "Gestion des Rejets : Rejet automatique des CDRs si le crédit est insuffisant ou l'abonné inconnu."
        ]
    },
    {
        "title": "Partie 2 : Interfaces et Fonctionnalités",
        "bullets": [
            "Simulateur de CDRs (index.php) : Injection d'événements réseau et affichage en temps réel.",
            "Module de Rechargement (recharger.php) : Crédit des comptes prépayés.",
            "Module d'Impression (facture.php) : Édition de reçus individuels.",
            "Exportation Comptable (exporter_csv.php) : Extraction globale vers Microsoft Excel."
        ]
    },
    {
        "title": "Conclusion et Bilan",
        "bullets": [
            "Système réparti Client-Serveur 100% fonctionnel sous XAMPP (PHP / MySQL).",
            "Gestion réussie de la convergence des services télécoms.",
            "Automatisation complète : Ingestion -> Évaluation -> Débit -> Reporting.",
            "Merci pour votre attention ! (Place aux questions)"
        ]
    }
]

for slide_info in slides_data:
    if slide_info.get("is_title_slide"):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = slide_info["title"]
        subtitle.text = slide_info["subtitle"]
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = slide_info["title"]
        tf = body.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(slide_info["bullets"]):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = bullet
            p.level = 0

prs.save("Presentation_Facturation_Larissa.pptx")
print("Le fichier 'Presentation_Facturation_Larissa.pptx' a été créé avec succès !")